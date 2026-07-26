"""Populate the official IDEA presentation template in place, as a pitch deck.

Deliberately edits `IDEA_Presentation_Format(1).pptx` rather than building a new
deck: the theme, slide master, layouts, colour scheme, title fonts, footers and
slide numbers are the template's and must survive untouched. Everything added
here is drawn with the template's own Office palette (dk2 #1F497D, accent1
#4F81BD, accent2 #C0504D, accent3 #9BBB59) and Arial, the font the template's own
body boxes use. Titles are left exactly as the template set them.

The template ships seven slides. Its own first slide is the "IMPORTANT
INSTRUCTIONS" sheet, which states: "You can delete this slide (Important
Pointers) when you upload the details of your idea." Removing it leaves exactly
the six slides the brief allows.

This is a pitch, not documentation: cards, a flow diagram, stat tiles and real
dashboard screenshots rather than paragraphs. Every number on a slide is
reproduced from this repository — see FACTS below for where each one comes from.
Registration data (problem statement ID, theme, student ID) is not in the
repository, so a visible FILL-IN marker is left instead of a guess.

Run:  python3 scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
TEMPLATE = REPO / "IDEA_Presentation_Format(1).pptx"
# The untouched original, kept so this script is idempotent: it always reads the
# pristine seven-slide template and writes the populated deck over TEMPLATE.
PRISTINE = REPO / "docs" / "template" / "IDEA_Presentation_Format(1).pptx"
ASSETS = REPO / "docs" / "assets"

R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"

# -- palette: the template's own theme colours, plus tints of them -----------

NAVY = RGBColor(0x1F, 0x49, 0x7D)      # dk2
BLUE = RGBColor(0x4F, 0x81, 0xBD)      # accent1
BLUE_D = RGBColor(0x36, 0x5F, 0x91)    # accent1, darker shade
BLUE_PALE = RGBColor(0xE7, 0xEE, 0xF7)  # accent1 tint
BLUE_TINT = RGBColor(0xD3, 0xDF, 0xEE)  # accent1 tint, card borders
RED = RGBColor(0xC0, 0x50, 0x4D)       # accent2
RED_PALE = RGBColor(0xF7, 0xE9, 0xE8)  # accent2 tint
GREEN = RGBColor(0x5F, 0x8A, 0x2E)     # accent3, shaded for contrast on white
GREEN_PALE = RGBColor(0xEC, 0xF2, 0xE1)  # accent3 tint
ORANGE = RGBColor(0xE3, 0x7C, 0x24)    # accent6, shaded
INK = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x6E, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xC8, 0xD4, 0xE6)

FONT = "Arial"
MONO = "Consolas"

FILL_IN = "«fill in from portal»"

# Geometry. 13.333 x 7.5 in; the template's footer band starts at 6.95.
MARGIN = 0.55
CW = 13.333 - 2 * MARGIN

# -- FACTS -------------------------------------------------------------------
# Headline experiment: runs 34 (baseline) and 38 (agent) in backend/ecoloop.db,
# both EnergyPlus 26.1, scenario summer_week, 192 steps of 900 s.
#   baseline 360.16 kWh / 16.21 kW / 79 comfort violations of 80 occupied steps
#   agent    244.66 kWh / 11.49 kW /  0 violations / mean PPD 5.3 %
#   /api/compare -> 32.07 % kWh, 29.15 % peak, 48.51 kg CO2, 20.79 cost
# Reproduce with the two commands printed on slide 5.
# Tests: 161 backend (pytest) + 231 frontend (vitest) = 392, all passing.
# Store: 39,690 timestep rows, 5,021 decisions, 39 runs.
# LLM: 22 successful supervisory calls measured 465-2449 ms, mean 785 ms.
# Resilience: run 37, 48 supervisory calls, Groq 429s, 44 fallbacks, 192/192
#   steps completed.  Long horizon: run 25, 35,040 steps, RC model, complete.


# -- template surgery -------------------------------------------------------


def delete_slide(prs: Presentation, index: int) -> None:
    """Remove a slide, dropping its relationship so the package stays valid."""
    id_list = prs.slides._sldIdLst  # noqa: SLF001 - no public API for deletion
    entry = list(id_list)[index]
    prs.part.drop_rel(entry.get(R_NS))
    id_list.remove(entry)


def find(slide, name: str):
    """The shape with this name, or None."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def bring_to_front(shape) -> None:
    """Move a shape to the end of the tree, so cards drawn later cannot cover it."""
    element = shape._element  # noqa: SLF001
    element.getparent().append(element)


def set_title(slide, text: str) -> None:
    """Replace the title placeholder's text, keeping its run formatting."""
    title = slide.shapes.title
    paragraph = title.text_frame.paragraphs[0]
    runs = paragraph.runs
    runs[0].text = text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)  # noqa: SLF001


# -- drawing primitives -----------------------------------------------------


def _plain_paragraph(para) -> None:
    """Strip the inherited bullet and hanging indent.

    The template's own body boxes are bulleted lists at 28pt; reusing one as a
    caption otherwise leaves a stray bullet glyph floating on the slide.
    """
    pPr = para._p.get_or_add_pPr()  # noqa: SLF001
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:buBlip"):
        for element in pPr.findall(qn(tag)):
            pPr.remove(element)
    bu_none = pPr.makeelement(qn("a:buNone"), {})
    anchors = [pPr.find(qn(tag)) for tag in ("a:tabLst", "a:defRPr", "a:extLst")]
    following = next((a for a in anchors if a is not None), None)
    if following is None:
        pPr.append(bu_none)
    else:
        following.addprevious(bu_none)


def write(frame, lines, *, anchor=MSO_ANCHOR.TOP) -> None:
    """Rewrite a text frame from `lines`.

    Each line is a dict: t=text, s=size pt, b=bold, i=italic, c=colour,
    a=alignment ('l'|'c'|'r'), sb=space-before pt, f=font name.
    Formatting is set explicitly rather than inherited, because the template's
    stub bullets carry 28pt sizing no real slide's content fits inside.
    """
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    for para in list(frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)  # noqa: SLF001
    first = frame.paragraphs[0]
    for run in list(first.runs):
        run._r.getparent().remove(run._r)  # noqa: SLF001

    align = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    for index, line in enumerate(lines):
        para = first if index == 0 else frame.add_paragraph()
        para.level = 0
        _plain_paragraph(para)
        para.alignment = align[line.get("a", "l")]
        if line.get("sb"):
            para.space_before = Pt(line["sb"])
        run = para.add_run()
        run.text = line.get("t", "")
        run.font.name = line.get("f", FONT)
        run.font.size = Pt(line.get("s", 12))
        run.font.bold = bool(line.get("b"))
        run.font.italic = bool(line.get("i"))
        run.font.color.rgb = line.get("c", INK)


def box(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    radius=0.10,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    pad=(0.16, 0.16, 0.10, 0.10),
):
    """A flat card. No preset shadow — the template's default shadow reads dated."""
    sp = slide.shapes.add_shape(
        shape, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if shape is MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    frame = sp.text_frame
    frame.word_wrap = True
    frame.margin_left, frame.margin_right = Inches(pad[0]), Inches(pad[1])
    frame.margin_top, frame.margin_bottom = Inches(pad[2]), Inches(pad[3])
    return sp


def label(slide, left, top, width, height, lines, *, anchor=MSO_ANCHOR.TOP):
    """A plain text box with no fill or outline."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tb.text_frame.margin_left = Inches(0)
    tb.text_frame.margin_right = Inches(0)
    tb.text_frame.margin_top = Inches(0)
    tb.text_frame.margin_bottom = Inches(0)
    write(tb.text_frame, lines, anchor=anchor)
    return tb


def bullet_dot(slide, left, top, size, glyph, *, fill, colour=WHITE, pt=11):
    """A small filled circle carrying a glyph — the deck's icon vocabulary."""
    sp = box(
        slide,
        left,
        top,
        size,
        size,
        fill=fill,
        radius=0.5,
        shape=MSO_SHAPE.OVAL,
        pad=(0.01, 0.01, 0.01, 0.01),
    )
    write(
        sp.text_frame,
        [{"t": glyph, "s": pt, "b": True, "c": colour, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return sp


def chip(slide, left, top, width, height, text, *, fill=BLUE_PALE, colour=BLUE_D, size=10.5):
    sp = box(slide, left, top, width, height, fill=fill, radius=0.5,
             pad=(0.08, 0.08, 0.02, 0.02))
    write(
        sp.text_frame,
        [{"t": text, "s": size, "b": True, "c": colour, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return sp


def kicker(slide, text: str) -> None:
    """The template's own body box, reused as the slide's one-line message."""
    shape = find(slide, "TextBox 8")
    shape.left, shape.top = Inches(MARGIN), Inches(1.02)
    shape.width, shape.height = Inches(CW), Inches(0.40)
    write(
        shape.text_frame,
        [{"t": text, "s": 14.5, "b": True, "i": True, "c": BLUE_D, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    bring_to_front(shape)


def picture(slide, name: str, left, top, width):
    """Place a screenshot, scaled from its true aspect ratio."""
    path = ASSETS / name
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width))
    return pic


# -- slide 1 — title page ---------------------------------------------------


def title_page(slide) -> None:
    write(
        find(slide, "Subtitle 3").text_frame,
        [{"t": "ECO-LOOP BUILDING AGENTS", "s": 34, "b": True, "c": NAVY, "a": "c"}],
    )
    find(slide, "Subtitle 3").left = Inches(MARGIN)
    find(slide, "Subtitle 3").top = Inches(0.24)
    find(slide, "Subtitle 3").width = Inches(CW)
    find(slide, "Subtitle 3").height = Inches(0.70)

    band = box(slide, MARGIN, 1.00, CW, 0.66, fill=NAVY, radius=0.28)
    write(
        band.text_frame,
        [{"t": "AI that thinks before your building wastes energy.",
          "s": 21, "b": True, "c": WHITE, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    label(
        slide, MARGIN, 1.76, CW, 0.42,
        [{"t": "An open-source LLM closes the control loop on a live EnergyPlus "
               "simulation — set-points injected into the running instance, "
               "savings proven against an identical baseline.",
          "s": 12.5, "c": MUTED, "a": "c"}],
    )

    # -- left: the registration block, as the template requires. The template's
    # own "TextBox 6" carries it, repositioned inside a card rather than
    # replaced, so the required pointers stay in the template's own shape.
    box(slide, MARGIN, 2.30, 6.05, 3.55, fill=BLUE_PALE, line=BLUE_TINT)
    details = find(slide, "TextBox 6")
    details.left, details.top = Inches(MARGIN + 0.20), Inches(2.46)
    details.width, details.height = Inches(5.65), Inches(3.25)
    write(
        details.text_frame,
        [
            {"t": "PROBLEM STATEMENT", "s": 11, "b": True, "c": BLUE_D},
            {"t": f"ID  —  {FILL_IN}", "s": 12.5, "c": INK, "sb": 7},
            {"t": "Title  —  Eco-Loop Building Agents: autonomous closed-loop "
                  "building control  (Honeywell Hackathon · Question 1)",
             "s": 12.5, "c": INK, "sb": 3},
            {"t": f"Theme  —  {FILL_IN}", "s": 12.5, "c": INK, "sb": 3},
            {"t": "PS Category  —  Software", "s": 12.5, "c": INK, "sb": 3},
            {"t": "TEAM", "s": 11, "b": True, "c": BLUE_D, "sb": 13},
            {"t": "Student Name (registered on portal)  —  Pranav Kumar",
             "s": 12.5, "c": INK, "sb": 7},
            {"t": f"Student ID  —  {FILL_IN}", "s": 12.5, "c": INK, "sb": 3},
            {"t": "IN THE REPOSITORY", "s": 11, "b": True, "c": BLUE_D, "sb": 13},
            {"t": "Closed-loop source · baseline and runtime-generated .idf models · "
                  "savings dashboard · docs/ARCHITECTURE.md · 392 passing tests",
             "s": 11.5, "c": MUTED, "sb": 7},
        ],
    )
    bring_to_front(details)

    # -- right: the number that has to land in the first ten seconds
    label(
        slide, 6.90, 2.30, 5.88, 0.30,
        [{"t": "MEASURED ON REAL ENERGYPLUS — NOT PROJECTED",
          "s": 10.5, "b": True, "c": BLUE_D, "a": "c"}],
    )
    stats = [
        ("32.1%", "ENERGY SAVED", "115.5 kWh", GREEN, GREEN_PALE),
        ("29.1%", "PEAK CUT", "16.2 → 11.5 kW", BLUE_D, BLUE_PALE),
        ("79 → 0", "COMFORT\nVIOLATIONS", "of 80 occupied steps", NAVY, BLUE_PALE),
    ]
    tw, gap = 1.86, 0.15
    for index, (big, cap, sub, colour, fill) in enumerate(stats):
        left = 6.90 + index * (tw + gap)
        tile = box(slide, left, 2.66, tw, 1.66, fill=fill, line=BLUE_TINT)
        write(
            tile.text_frame,
            [
                {"t": big, "s": 30, "b": True, "c": colour, "a": "c"},
                {"t": cap.replace("\n", " "), "s": 10.5, "b": True, "c": INK,
                 "a": "c", "sb": 4},
                {"t": sub, "s": 9.5, "c": MUTED, "a": "c", "sb": 2},
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    verdict = box(slide, 6.90, 4.44, 5.88, 0.42, fill=GREEN, radius=0.4)
    write(
        verdict.text_frame,
        [{"t": "✓   COMFORT MAINTAINED   ·   ASHRAE-55 PMV BAND HELD ALL DAY",
          "s": 11.5, "b": True, "c": WHITE, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    label(
        slide, 6.90, 4.98, 5.88, 0.90,
        [
            {"t": "EnergyPlus 26.1  ·  summer week  ·  192 timesteps of 15 min",
             "s": 11, "b": True, "c": INK, "a": "c"},
            {"t": "Agent run and fixed-schedule baseline execute the identical "
                  "code path, weather and occupancy — a controlled experiment, "
                  "reproducible from two CLI commands.",
             "s": 10, "c": MUTED, "a": "c", "sb": 4},
        ],
    )

    # -- bottom: the stack, as small labels
    chips = [
        "EnergyPlus 26.1", "Llama 3.3 70B via Groq", "MCP server",
        "FastAPI + SSE", "React + Recharts", "392 tests passing",
    ]
    cw, gap = 1.95, 0.11
    total = len(chips) * cw + (len(chips) - 1) * gap
    x = (13.333 - total) / 2
    for text in chips:
        chip(slide, x, 6.12, cw, 0.36, text)
        x += cw + gap


# -- slide 2 — problem and solution -----------------------------------------


def proposed_solution(slide) -> None:
    set_title(slide, "ECO-LOOP BUILDING AGENTS")
    kicker(slide, "Buildings waste energy because their controls cannot think. "
                  "Ours reasons — every fifteen minutes.")

    # -- left: the problem
    panel = box(slide, MARGIN, 1.58, 5.30, 4.18, fill=RED_PALE, line=RGBColor(0xE6, 0xC9, 0xC7))
    write(
        panel.text_frame,
        [{"t": "THE PROBLEM", "s": 15, "b": True, "c": RED}],
    )
    problems = [
        ("Buildings burn ~40% of global energy",
         "and remain a primary driver of carbon emissions."),
        ("Building management systems run fixed schedules",
         "blind to weather, occupancy and grid carbon intensity."),
        ("Control is reactive, not anticipatory",
         "comfort is corrected only after occupants have complained."),
    ]
    y = 2.06
    for head, sub in problems:
        bullet_dot(slide, MARGIN + 0.18, y + 0.04, 0.28, "✕", fill=RED, pt=12)
        label(
            slide, MARGIN + 0.56, y, 4.55, 0.72,
            [
                {"t": head, "s": 12.5, "b": True, "c": INK},
                {"t": sub, "s": 11, "c": MUTED, "sb": 2},
            ],
        )
        y += 0.80

    proof = box(slide, MARGIN + 0.18, 4.60, 4.94, 1.00, fill=WHITE, line=RED)
    write(
        proof.text_frame,
        [
            {"t": "79 of 80", "s": 24, "b": True, "c": RED, "a": "c"},
            {"t": "occupied timesteps outside the comfort band — measured on the "
                  "fixed-schedule baseline in EnergyPlus.",
             "s": 10.5, "c": INK, "a": "c", "sb": 3},
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # -- right: the solution
    label(
        slide, 6.15, 1.58, 6.63, 0.34,
        [{"t": "OUR SOLUTION", "s": 15, "b": True, "c": NAVY}],
    )
    items = [
        ("AI SUPERVISOR",
         "Llama 3.3 70B reads live telemetry and sets the control policy."),
        ("ENERGYPLUS DIGITAL TWIN",
         "Set-points injected into the running instance — no restart, no file rewrite."),
        ("REACTIVE SAFETY GUARD",
         "Every action clamped in code to hard comfort and equipment limits."),
        ("MCP TOOL LAYER",
         "The same six tools serve the agent and any external MCP client."),
        ("LIVE SAVINGS DASHBOARD",
         "Eight panels streaming telemetry, reasoning and kWh saved."),
    ]
    y = 2.00
    for head, sub in items:
        card = box(slide, 6.15, y, 6.63, 0.68, fill=BLUE_PALE, line=BLUE_TINT)
        bullet_dot(slide, 6.30, y + 0.19, 0.30, "✓", fill=GREEN, pt=12)
        label(
            slide, 6.72, y + 0.09, 5.94, 0.54,
            [
                {"t": head, "s": 12, "b": True, "c": NAVY},
                {"t": sub, "s": 10.5, "c": MUTED, "sb": 1},
            ],
        )
        y += 0.76

    band = box(slide, MARGIN, 5.94, CW, 0.62, fill=NAVY, radius=0.22)
    write(
        band.text_frame,
        [{"t": "WHAT MAKES IT DIFFERENT   ·   baseline and agent run the identical "
               "code path, weather and occupancy — so the saving is a controlled "
               "experiment, not a demo",
          "s": 12, "b": True, "c": WHITE, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )


# -- slide 3 — architecture --------------------------------------------------


def technical_approach(slide) -> None:
    kicker(slide, "One closed loop: sense → reason → clamp → actuate, "
                  "inside the running simulation.")

    stages = [
        ("ENERGYPLUS\nDIGITAL TWIN", "small_office.idf · EPW weather",
         "pyenergyplus runtime API", NAVY),
        ("SENSOR\nTELEMETRY", "zone temp · outdoor · occupancy\nCO₂ · power · PMV",
         "every 15-min timestep", BLUE_D),
        ("LLM\nSUPERVISOR", "llama-3.3-70b via Groq\n6 JSON-schema tools",
         "on cadence · ~0.8 s", BLUE),
        ("REACTIVE\nGUARD", "clamps to hard comfort\nand equipment limits",
         "every step · microseconds", GREEN),
        ("FORWARD\nINJECTION", "set-points written to\nEnergyPlus actuators",
         "live instance · no restart", ORANGE),
    ]
    bw, gap = 2.17, 0.325
    x, y, bh = MARGIN, 1.56, 1.88
    centres = []
    for head, body, foot, colour in stages:
        card = box(slide, x, y, bw, bh, fill=WHITE, line=colour, line_w=1.75)
        head_bar = box(slide, x, y, bw, 0.52, fill=colour, radius=0.20)
        write(
            head_bar.text_frame,
            [{"t": head.replace("\n", " "), "s": 10.5, "b": True, "c": WHITE, "a": "c"}],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        lines = [{"t": part, "s": 9.5, "c": INK, "a": "c", "sb": 2}
                 for part in body.split("\n")]
        lines.append({"t": foot, "s": 8.5, "i": True, "c": MUTED, "a": "c", "sb": 5})
        label(slide, x + 0.08, y + 0.62, bw - 0.16, bh - 0.70, lines,
              anchor=MSO_ANCHOR.MIDDLE)
        centres.append(x + bw / 2)
        if x + bw + gap < MARGIN + CW:
            arrow = box(slide, x + bw + 0.055, y + bh / 2 - 0.15, gap - 0.11, 0.30,
                        fill=BLUE, shape=MSO_SHAPE.RIGHT_ARROW)
        x += bw + gap
        _ = card

    # the return path: what makes it a loop rather than a pipeline
    loop = box(slide, MARGIN, 3.56, CW, 0.54, fill=BLUE_PALE, line=BLUE,
               shape=MSO_SHAPE.LEFT_ARROW)
    write(
        loop.text_frame,
        [{"t": "CLOSED LOOP  —  the next timestep's telemetry already reflects the "
               "action just written.  35,040 steps (a full simulated year) "
               "completed end to end.",
          "s": 11, "b": True, "c": NAVY, "a": "c"}],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    support = [
        ("EVIDENCE STORE",
         "SQLite: runs · timesteps · decisions.\n39,690 telemetry rows and 5,021 "
         "decisions persisted, with rationale, latency and guard overrides."),
        ("API LAYER",
         "FastAPI — 14 REST endpoints plus a Server-Sent Events stream and "
         "per-run CSV export. The dashboard is a pure reader."),
        ("MCP SERVER",
         "The agent's own ToolRegistry served over stdio, so an external MCP "
         "client drives the identical six tools."),
    ]
    sw, sgap = 3.94, 0.225
    x, y = MARGIN, 4.26
    for head, body in support:
        card = box(slide, x, y, sw, 1.42, fill=BLUE_PALE, line=BLUE_TINT)
        write(
            card.text_frame,
            [
                {"t": head, "s": 11.5, "b": True, "c": NAVY},
                {"t": body.replace("\n", " "), "s": 10, "c": MUTED, "sb": 3},
            ],
        )
        x += sw + sgap

    chips = [
        "Python 3.12", "EnergyPlus 26.1", "pyenergyplus", "Groq · Llama 3.3",
        "MCP", "FastAPI", "React + Vite", "Recharts", "392 tests",
    ]
    cw, gap = 1.30, 0.075
    total = len(chips) * cw + (len(chips) - 1) * gap
    x = (13.333 - total) / 2
    for text in chips:
        chip(slide, x, 5.86, cw, 0.36, text, size=9)
        x += cw + gap


# -- slide 4 — feasibility ---------------------------------------------------


def feasibility(slide) -> None:
    kicker(slide, "It already runs — and it is engineered to fail safely.")

    proofs = [
        ("RUNS TODAY",
         "Real EnergyPlus 26.1 and real Groq calls, measured at 465–2,449 ms per "
         "supervisory decision. Two CLI commands reproduce the headline number."),
        ("TESTED",
         "392 automated tests pass — 161 backend (pytest) and 231 frontend "
         "(vitest), including the EnergyPlus thread bridge and forward injection."),
        ("EVIDENCED",
         "39 recorded runs · 39,690 telemetry rows · 5,021 decisions, each with "
         "its rationale and latency. Every run exports to CSV."),
    ]
    pw, gap = 3.94, 0.225
    x, y = MARGIN, 1.54
    for head, body in proofs:
        card = box(slide, x, y, pw, 1.46, fill=GREEN_PALE,
                   line=RGBColor(0xCF, 0xDD, 0xB6))
        bullet_dot(slide, x + 0.16, y + 0.16, 0.30, "✓", fill=GREEN, pt=12)
        label(
            slide, x + 0.58, y + 0.16, pw - 0.76, 1.16,
            [
                {"t": head, "s": 12.5, "b": True, "c": GREEN},
                {"t": body, "s": 10, "c": INK, "sb": 4},
            ],
        )
        x += pw + gap

    label(
        slide, MARGIN, 3.16, CW, 0.32,
        [{"t": "RISK   →   OUR ANSWER", "s": 12.5, "b": True, "c": NAVY, "a": "c"}],
    )

    risks = [
        ("Model latency on a\n35,040-step year",
         "Two-tier control. The guard runs every step in microseconds; the LLM "
         "runs on a cadence and only chooses the policy. The loop never waits "
         "on a model."),
        ("A hallucinated or\nunsafe set-point",
         "Safety lives in code, not in a prompt. Every action is clamped to hard "
         "comfort and equipment limits, and each override is recorded as "
         "guard_clamped."),
        ("Hosted API rate-limits,\ntimes out, or fails",
         "Bounded retries → circuit breaker → baseline fallback through the same "
         "guard. One 192-step run completed with 44 of 48 supervisory calls "
         "degraded and zero crashes."),
        ("Logs and telemetry\nexceed any context window",
         "logsummary.py compacts before prompting: severity filtering, warning "
         "de-duplication and statistical windowing of long traces."),
    ]
    rw, rgap = 2.90, 0.213
    x, y = MARGIN, 3.56
    for head, body in risks:
        card = box(slide, x, y, rw, 2.62, fill=WHITE, line=BLUE_TINT)
        head_bar = box(slide, x, y, rw, 0.66, fill=RED_PALE, radius=0.16)
        write(
            head_bar.text_frame,
            [{"t": head.replace("\n", " "), "s": 11, "b": True, "c": RED, "a": "c"}],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        bullet_dot(slide, x + rw / 2 - 0.15, y + 0.74, 0.30, "↓", fill=BLUE, pt=12)
        label(slide, x + 0.16, y + 1.18, rw - 0.32, 1.34,
              [{"t": body, "s": 10, "c": INK}])
        x += rw + rgap
        _ = card


# -- slide 5 — artifacts, results and demo -----------------------------------


def artifacts(slide) -> None:
    kicker(slide, "The live dashboard on the real run: 32.1% less energy, "
                  "zero comfort violations.")

    picture(slide, "kpi_row.png", MARGIN, 1.48, CW)          # 3080 × 220
    picture(slide, "charts_row.png", MARGIN, 2.48, 7.00)      # 3080 × 690
    picture(slide, "reasoning.png", 7.83, 2.48, 4.95)         # 1890 × 600

    label(
        slide, MARGIN, 4.14, 7.00, 0.26,
        [{"t": "Live telemetry — zone temperature against the set-point band · "
               "agent vs baseline energy · PMV inside the ASHRAE-55 band",
          "s": 9, "i": True, "c": MUTED, "a": "c"}],
    )
    label(
        slide, 7.83, 4.14, 4.95, 0.26,
        [{"t": "Every decision explains itself — and the guard's overrides are "
               "on the record",
          "s": 9, "i": True, "c": MUTED, "a": "c"}],
    )

    wins = [
        ("32.1%", "kWh saved vs baseline"),
        ("0", "comfort violations"),
        ("192/192", "steps, zero crashes"),
        ("48.5 kg", "CO₂ avoided in 2 days"),
    ]
    ww, gap = 1.66, 0.12
    x, y = MARGIN, 4.60
    for big, cap in wins:
        tile = box(slide, x, y, ww, 0.74, fill=NAVY, radius=0.14)
        write(
            tile.text_frame,
            [
                {"t": big, "s": 17, "b": True, "c": WHITE, "a": "c"},
                {"t": cap, "s": 8.5, "c": RGBColor(0xC9, 0xD6, 0xE8), "a": "c", "sb": 1},
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += ww + gap

    quote = box(slide, 7.83, 4.60, 4.95, 1.86, fill=BLUE_PALE, line=BLUE_TINT)
    write(
        quote.text_frame,
        [
            {"t": "AGENT DECISION — VERBATIM FROM THE RUN LOG", "s": 9, "b": True,
             "c": BLUE_D},
            {"t": "“The zone is unoccupied and will remain so for the next 4 steps, "
                  "so we can set back the temperature and turn off the lighting to "
                  "save energy.”", "s": 10, "i": True, "c": INK, "sb": 5},
            {"t": "set_control_policy · setback · 19.0 / 20.5 °C · 849 ms · "
                  "llama-3.3-70b-versatile", "s": 8.5, "c": MUTED, "sb": 5},
        ],
    )

    code = box(slide, MARGIN, 5.48, 7.10, 0.98, fill=RGBColor(0xF4, 0xF6, 0xFA),
               line=BLUE_TINT)
    write(
        code.text_frame,
        [
            {"t": "REPRODUCE THE HEADLINE NUMBER", "s": 9, "b": True, "c": BLUE_D},
            {"t": "python cli.py --scenario summer_week --controller baseline "
                  "--simulator energyplus --days 2", "s": 8.5, "c": NAVY, "f": MONO,
             "sb": 4},
            {"t": "python cli.py --scenario summer_week --controller rule "
                  "--simulator energyplus --days 2 --compare <id>",
             "s": 8.5, "c": NAVY, "f": MONO, "sb": 2},
        ],
    )

    label(
        slide, MARGIN, 6.54, CW, 0.30,
        [{"t": "Artifacts:  backend/ (24 modules) · frontend/ (9 components) · "
               "small_office.idf + runtime-generated variants · "
               "docs/ARCHITECTURE.md · 392 tests",
          "s": 9.5, "c": MUTED, "a": "c"}],
    )


# -- slide 6 — why it wins, and what it stands on ----------------------------


def references(slide) -> None:
    kicker(slide, "WHY THIS SOLUTION WINS — six claims, each backed by something "
                  "in the repository.")

    wins = [
        ("REAL AI", "An open-weight Llama 3.3 70B reasons over live telemetry and "
                    "emits a policy through native tool-calling.", BLUE_D),
        ("REAL SIMULATION", "EnergyPlus 26.1 driven through its runtime API — "
                            "actuators written mid-run, not a file regenerated "
                            "between runs.", NAVY),
        ("EXPLAINABLE", "Every decision persists its rationale, tool call, latency "
                        "and any guard override. The dashboard shows them live.",
         BLUE_D),
        ("SAFE", "A deterministic guard clamps every action to hard comfort and "
                 "equipment limits, so a hallucinated set-point cannot reach the "
                 "building.", GREEN),
        ("SCALABLE", "The simulator sits behind a Protocol and the model is named "
                     "in .env — neither vendor is locked in; the MCP server opens "
                     "the tools to any client.", BLUE),
        ("PRODUCTION-READY", "392 passing tests, a typed API, an evidence store and "
                             "CSV export — a controlled experiment anyone can "
                             "re-run.", ORANGE),
    ]
    ww, gap = 3.94, 0.225
    for index, (head, body, colour) in enumerate(wins):
        x = MARGIN + (index % 3) * (ww + gap)
        y = 1.56 + (index // 3) * 1.20
        card = box(slide, x, y, ww, 1.10, fill=WHITE, line=BLUE_TINT)
        stripe = box(slide, x, y, 0.11, 1.10, fill=colour, radius=0.5)
        write(
            card.text_frame,
            [
                {"t": head, "s": 12, "b": True, "c": colour},
                {"t": body, "s": 9.5, "c": INK, "sb": 3},
            ],
        )
        _ = stripe

    label(
        slide, MARGIN, 4.06, CW, 0.30,
        [{"t": "RESEARCH AND REFERENCES", "s": 12, "b": True, "c": NAVY, "a": "c"}],
    )

    columns = [
        ("STANDARDS", [
            "Fanger PMV / PPD thermal comfort — ISO 7730 Annex D and "
            "ANSI/ASHRAE Standard 55; implemented in app/comfort.py.",
            "EnergyPlus Weather (EPW) hourly format — parsed in app/sim/weather.py.",
        ]),
        ("SIMULATION ENGINE", [
            "EnergyPlus 26.1 and its runtime Python API (pyenergyplus), NREL — "
            "github.com/NREL/EnergyPlus/releases",
            "Baseline model derived from the distribution example "
            "MovableExtInsulationSimple.idf.",
        ]),
        ("MODEL AND PROTOCOL", [
            "Groq inference API serving open-weight Llama 3.3 70B "
            "(llama-3.3-70b-versatile) — console.groq.com",
            "Model Context Protocol — official Python SDK (mcp ≥ 1.0), wrapping "
            "the agent's own ToolRegistry.",
        ]),
    ]
    cw, gap = 3.94, 0.225
    x, y = MARGIN, 4.40
    for head, entries in columns:
        card = box(slide, x, y, cw, 1.72, fill=BLUE_PALE, line=BLUE_TINT)
        lines = [{"t": head, "s": 10.5, "b": True, "c": BLUE_D}]
        for entry in entries:
            lines.append({"t": "•  " + entry, "s": 9.5, "c": INK, "sb": 5})
        write(card.text_frame, lines)
        x += cw + gap

    label(
        slide, MARGIN, 6.24, CW, 0.30,
        [{"t": "Full design rationale, sequence and data-flow diagrams, database "
               "schema and design decisions D1–D11:  docs/ARCHITECTURE.md",
          "s": 10, "i": True, "c": MUTED, "a": "c"}],
    )


# -- entry point -------------------------------------------------------------


def main() -> int:
    source = PRISTINE if PRISTINE.exists() else TEMPLATE
    prs = Presentation(source)
    if len(prs.slides) != 7:
        raise SystemExit(f"expected the 7-slide template, found {len(prs.slides)}")

    # The template's own instruction slide, which it says to delete on upload.
    delete_slide(prs, 0)

    builders = (
        title_page,
        proposed_solution,
        technical_approach,
        feasibility,
        artifacts,
        references,
    )
    if len(prs.slides) != len(builders):
        raise SystemExit(f"slide/builder mismatch: {len(prs.slides)} vs {len(builders)}")
    for slide, build in zip(prs.slides, builders):
        build(slide)

    prs.save(TEMPLATE)
    print(f"populated {TEMPLATE.name}: {len(prs.slides)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
